import streamlit as st
from pptx import Presentation
from PyPDF2 import PdfReader

def extract_text_from_pptx(file):
    """Extract text from a PowerPoint file."""
    presentation = Presentation(file)
    text_data = []
    for i, slide in enumerate(presentation.slides):
        slide_text = f"Slide {i + 1}:\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"
        text_data.append(slide_text)
    return "\n".join(text_data)

def extract_text_from_pdf(file):
    """Extract text from a PDF file."""
    reader = PdfReader(file)
    text_data = []
    for page in reader.pages:
        text_data.append(page.extract_text())
    return "\n".join(text_data)

# Streamlit UI
st.title("Ekstraktor Teks untuk File PowerPoint dan PDF")
st.write("Dibuat oleh: Harits Raharjo Setiono")

uploaded_file = st.file_uploader("Unggah file .pptx atau .pdf", type=["pptx", "pdf"])

if uploaded_file is not None:
  file_extension = uploaded_file.name.split(".")[-1].lower()
  
  if file_extension == "pptx":
    st.write("Mengekstrak teks dari PowerPoint...")
    text = extract_text_from_pptx(uploaded_file)
    st.text_area("Teks yang Diekstrak:", text, height=400)
  elif file_extension == "pdf":
    st.write("Mengekstrak teks dari PDF...")
    text = extract_text_from_pdf(uploaded_file)
    st.text_area("Teks yang Diekstrak:", text, height=400)
  else:
    st.error("Jenis file tidak didukung. Silakan unggah file .pptx atau .pdf.")

st.write("Versi 1.0")